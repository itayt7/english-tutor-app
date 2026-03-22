"""
Document parsing service for PDF and PPTX files.

Uses PyMuPDF (fitz) for PDFs and python-pptx for PowerPoint files.
Extracts text slide-by-slide / page-by-page and returns structured data.
"""

import io
import logging
from typing import List

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from app.schemas.presentation import ExtractedSlide, DocumentExtractionResult

logger = logging.getLogger(__name__)

# Supported MIME types mapped to their format
SUPPORTED_CONTENT_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
}

SUPPORTED_EXTENSIONS = {".pdf", ".pptx"}


def _extract_from_pdf(file_bytes: bytes) -> List[ExtractedSlide]:
    """Extract text page-by-page from a PDF using PyMuPDF."""
    slides: List[ExtractedSlide] = []
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
    except Exception as exc:
        logger.error(f"Failed to open PDF: {exc}")
        raise ValueError(f"Could not parse PDF file: {exc}") from exc

    for page_num in range(len(doc)):
        page = doc[page_num]
        text = page.get_text("text").strip()
        slides.append(
            ExtractedSlide(
                page_number=page_num + 1,  # 1-based
                text=text,
            )
        )

    doc.close()
    return slides


def _extract_from_pptx(file_bytes: bytes) -> List[ExtractedSlide]:
    """Extract text slide-by-slide from a PPTX using python-pptx."""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except PackageNotFoundError as exc:
        logger.error(f"Failed to open PPTX: {exc}")
        raise ValueError(f"Could not parse PPTX file: {exc}") from exc
    except Exception as exc:
        logger.error(f"Failed to open PPTX: {exc}")
        raise ValueError(f"Could not parse PPTX file: {exc}") from exc

    slides: List[ExtractedSlide] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        text_parts.append(paragraph_text)
        slides.append(
            ExtractedSlide(
                page_number=slide_num,
                text="\n".join(text_parts),
            )
        )

    return slides


def parse_document(filename: str, file_bytes: bytes) -> DocumentExtractionResult:
    """
    Parse a PDF or PPTX file and extract text from each page/slide.

    Args:
        filename: Original filename (used to determine format).
        file_bytes: Raw bytes of the uploaded file.

    Returns:
        DocumentExtractionResult with extracted slide data.

    Raises:
        ValueError: If the file format is unsupported or the file is corrupted.
    """
    lower_name = filename.lower()

    if lower_name.endswith(".pdf"):
        slides = _extract_from_pdf(file_bytes)
    elif lower_name.endswith(".pptx"):
        slides = _extract_from_pptx(file_bytes)
    else:
        raise ValueError(
            f"Unsupported file format: '{filename}'. "
            f"Only .pdf and .pptx files are supported."
        )

    return DocumentExtractionResult(filename=filename, slides=slides)
