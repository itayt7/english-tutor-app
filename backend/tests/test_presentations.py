"""
Tests for the Document Processing Engine (Task 1 — Phase 4).

Covers:
  • Schema validation
  • document_parser service (PDF + PPTX)
  • /api/v1/presentations/upload API endpoint
  • Edge cases: corrupted files, unsupported formats, scanned PDFs, empty files
"""

import io
import pytest
from unittest.mock import patch, MagicMock

from httpx import AsyncClient, ASGITransport

# ── Schemas ───────────────────────────────────────────────────────────────────
from app.schemas.presentation import ExtractedSlide, DocumentExtractionResult

# ── Service ───────────────────────────────────────────────────────────────────
from app.services.document_parser import (
    parse_document,
    _extract_from_pdf,
    _extract_from_pptx,
)

# ── FastAPI app (for integration tests) ──────────────────────────────────────
from main import app


# ===========================================================================
# 1. Schema unit tests
# ===========================================================================


class TestPresentationSchemas:
    """Validate the Pydantic models for document extraction."""

    def test_extracted_slide_valid(self):
        slide = ExtractedSlide(page_number=1, text="Hello world")
        assert slide.page_number == 1
        assert slide.text == "Hello world"

    def test_extracted_slide_page_number_must_be_positive(self):
        with pytest.raises(Exception):
            ExtractedSlide(page_number=0, text="Nope")

    def test_document_extraction_result_valid(self):
        result = DocumentExtractionResult(
            filename="deck.pptx",
            slides=[
                ExtractedSlide(page_number=1, text="Slide 1"),
                ExtractedSlide(page_number=2, text="Slide 2"),
            ],
        )
        assert result.filename == "deck.pptx"
        assert len(result.slides) == 2

    def test_document_extraction_result_empty_slides(self):
        result = DocumentExtractionResult(filename="empty.pdf", slides=[])
        assert result.slides == []


# ===========================================================================
# 2. document_parser service tests
# ===========================================================================


class TestDocumentParserPDF:
    """Tests for PDF extraction via PyMuPDF."""

    def _make_pdf_bytes(self, pages: list[str]) -> bytes:
        """Helper: create a real in-memory PDF with given page texts."""
        import fitz

        doc = fitz.open()
        for text in pages:
            page = doc.new_page()
            page.insert_text((72, 72), text)
        pdf_bytes = doc.tobytes()
        doc.close()
        return pdf_bytes

    def test_extract_single_page_pdf(self):
        pdf_bytes = self._make_pdf_bytes(["Page one content"])
        result = parse_document("test.pdf", pdf_bytes)
        assert result.filename == "test.pdf"
        assert len(result.slides) == 1
        assert result.slides[0].page_number == 1
        assert "Page one content" in result.slides[0].text

    def test_extract_multi_page_pdf(self):
        pages = [f"Content of page {i}" for i in range(1, 11)]
        pdf_bytes = self._make_pdf_bytes(pages)
        result = parse_document("deck.pdf", pdf_bytes)
        assert len(result.slides) == 10
        for i, slide in enumerate(result.slides, start=1):
            assert slide.page_number == i
            assert f"Content of page {i}" in slide.text

    def test_scanned_pdf_returns_empty_text(self):
        """A PDF with no extractable text should still return slides with empty text."""
        import fitz

        doc = fitz.open()
        doc.new_page()  # blank page, no text
        pdf_bytes = doc.tobytes()
        doc.close()

        result = parse_document("scanned.pdf", pdf_bytes)
        assert len(result.slides) == 1
        assert result.slides[0].text == ""

    def test_corrupted_pdf_raises_value_error(self):
        with pytest.raises(ValueError, match="Could not parse PDF"):
            parse_document("bad.pdf", b"this is not a pdf")


class TestDocumentParserPPTX:
    """Tests for PPTX extraction via python-pptx."""

    def _make_pptx_bytes(self, slide_texts: list[str]) -> bytes:
        """Helper: create a real in-memory PPTX with given slide texts."""
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        for text in slide_texts:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
            txBox.text_frame.text = text
        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    def test_extract_single_slide_pptx(self):
        pptx_bytes = self._make_pptx_bytes(["Slide one text"])
        result = parse_document("test.pptx", pptx_bytes)
        assert result.filename == "test.pptx"
        assert len(result.slides) == 1
        assert result.slides[0].page_number == 1
        assert "Slide one text" in result.slides[0].text

    def test_extract_multi_slide_pptx(self):
        texts = [f"Slide {i} content" for i in range(1, 11)]
        pptx_bytes = self._make_pptx_bytes(texts)
        result = parse_document("deck.pptx", pptx_bytes)
        assert len(result.slides) == 10
        for i, slide in enumerate(result.slides, start=1):
            assert slide.page_number == i
            assert f"Slide {i} content" in slide.text

    def test_pptx_with_empty_slides(self):
        """Slides with no text should still be returned."""
        from pptx import Presentation as PptxPresentation

        prs = PptxPresentation()
        prs.slides.add_slide(prs.slide_layouts[6])  # blank, no text
        buf = io.BytesIO()
        prs.save(buf)
        pptx_bytes = buf.getvalue()

        result = parse_document("empty_slides.pptx", pptx_bytes)
        assert len(result.slides) == 1
        assert result.slides[0].text == ""

    def test_corrupted_pptx_raises_value_error(self):
        with pytest.raises(ValueError, match="Could not parse PPTX"):
            parse_document("bad.pptx", b"this is not a pptx file")


class TestDocumentParserUnsupported:
    """Tests for unsupported file types."""

    def test_docx_rejected(self):
        with pytest.raises(ValueError, match="Unsupported file format"):
            parse_document("notes.docx", b"fake content")

    def test_txt_rejected(self):
        with pytest.raises(ValueError, match="Unsupported file format"):
            parse_document("readme.txt", b"some text")

    def test_no_extension_rejected(self):
        with pytest.raises(ValueError, match="Unsupported file format"):
            parse_document("noextension", b"data")


# ===========================================================================
# 3. API endpoint integration tests
# ===========================================================================


class TestPresentationsAPI:
    """Integration tests for POST /api/v1/presentations/upload."""

    async def _upload(self, file_bytes: bytes, filename: str, content_type: str = "application/octet-stream"):
        transport = ASGITransport(app=app)
        async with AsyncClient(transport=transport, base_url="http://test") as client:
            return await client.post(
                "/api/v1/presentations/upload",
                files={"file": (filename, file_bytes, content_type)},
            )

    # ── Happy paths ───────────────────────────────────────────────────────

    @pytest.mark.asyncio
    async def test_upload_valid_pdf(self):
        import fitz

        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((72, 72), "Hello from PDF")
        pdf_bytes = doc.tobytes()
        doc.close()

        resp = await self._upload(pdf_bytes, "test.pdf", "application/pdf")
        assert resp.status_code == 200
        body = resp.json()
        assert body["filename"] == "test.pdf"
        assert len(body["slides"]) == 1
        assert "Hello from PDF" in body["slides"][0]["text"]

    @pytest.mark.asyncio
    async def test_upload_valid_pptx(self):
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        txBox.text_frame.text = "Hello from PPTX"
        buf = io.BytesIO()
        prs.save(buf)

        resp = await self._upload(buf.getvalue(), "test.pptx")
        assert resp.status_code == 200
        body = resp.json()
        assert body["filename"] == "test.pptx"
        assert len(body["slides"]) == 1
        assert "Hello from PPTX" in body["slides"][0]["text"]

    @pytest.mark.asyncio
    async def test_upload_10_slide_deck_page_mapping(self):
        import fitz

        doc = fitz.open()
        for i in range(1, 11):
            page = doc.new_page()
            page.insert_text((72, 72), f"Page {i}")
        pdf_bytes = doc.tobytes()
        doc.close()

        resp = await self._upload(pdf_bytes, "big.pdf", "application/pdf")
        assert resp.status_code == 200
        body = resp.json()
        assert len(body["slides"]) == 10
        for i, slide in enumerate(body["slides"], start=1):
            assert slide["page_number"] == i
            assert f"Page {i}" in slide["text"]

    # ── Error paths ───────────────────────────────────────────────────────

    @pytest.mark.asyncio
    async def test_upload_unsupported_format_returns_415(self):
        resp = await self._upload(b"fake docx", "notes.docx")
        assert resp.status_code == 415

    @pytest.mark.asyncio
    async def test_upload_no_extension_returns_415(self):
        resp = await self._upload(b"data", "noext")
        assert resp.status_code == 415

    @pytest.mark.asyncio
    async def test_upload_corrupted_pdf_returns_422(self):
        resp = await self._upload(b"not a real pdf", "broken.pdf", "application/pdf")
        assert resp.status_code == 422

    @pytest.mark.asyncio
    async def test_upload_corrupted_pptx_returns_422(self):
        resp = await self._upload(b"not a real pptx", "broken.pptx")
        assert resp.status_code == 422

    @pytest.mark.asyncio
    async def test_upload_empty_file_returns_422(self):
        resp = await self._upload(b"", "empty.pdf", "application/pdf")
        assert resp.status_code == 422
